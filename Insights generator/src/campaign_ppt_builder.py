import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.visuals import latest_week_dimension_bar


# -----------------------------------------------------------------------------
# CHART COORDINATES (your template)
# -----------------------------------------------------------------------------
COORDS = {
    "Spend":      {"x": Cm(0.41),  "y": Cm(12.48), "w": Cm(7.5), "h": Cm(6)},
    "Traffic":    {"x": Cm(8.44),  "y": Cm(12.48), "w": Cm(7.5), "h": Cm(6)},
    "Engaged":    {"x": Cm(16.48), "y": Cm(12.48), "w": Cm(7.5), "h": Cm(6)},
    "EVR":        {"x": Cm(24.71), "y": Cm(12.48), "w": Cm(7.5), "h": Cm(6)},
}


# -----------------------------------------------------------------------------
# KPI BOX COORDINATES (top row)
# -----------------------------------------------------------------------------
KPI_COORDS = {
    "Spend":      {"x": Cm(0.6),  "y": Cm(2.2)},
    "Traffic":    {"x": Cm(8.6),  "y": Cm(2.2)},
    "Engaged":    {"x": Cm(16.6), "y": Cm(2.2)},
    "EVR":        {"x": Cm(24.7), "y": Cm(2.2)},
}

KPI_ICONS = {
    "Spend": "£",
    "Traffic": "👣",
    "Engaged": "⭐",
    "EVR": "📊",
}


# -----------------------------------------------------------------------------
# SUMMARY TEXT BLOCK
# -----------------------------------------------------------------------------
SUMMARY_COORDS = {"x": Cm(0.6), "y": Cm(6.8), "w": Cm(27), "h": Cm(4.0)}

TABLE_COORDS = {"x": Cm(0.6), "y": Cm(10.5), "w": Cm(12), "h": Cm(2.7)}



# ============================================================================
# → KPI FUNCTIONS
# ============================================================================

def format_kpi_value(key, v):
    if key == "Spend":
        return f"£{v:,.0f}"
    if key == "EVR":
        return f"{v*100:.1f}%"
    return f"{v:,.0f}"


def colour_for_delta(delta):
    """
    Returns RGB colors:
    ↑ green = positive movement
    ↓ red = negative movement
    → grey = stable / no baseline
    """

    if delta is None:
        return RGBColor(120, 120, 120)

    if delta > 0.02:
        return RGBColor(0, 150, 0)       # green
    if delta < -0.02:
        return RGBColor(200, 0, 0)       # red
    return RGBColor(120, 120, 120)       # neutral


def get_delta(curr, prev):
    if prev == 0 or prev is None:
        return None
    return (curr - prev) / prev


def extract_latest_kpis(df):
    df = df.sort_values(["Year", "Quarter", "Week"])
    current = df.tail(1).iloc[0]
    prev = df.tail(2).iloc[0] if len(df) >= 2 else None

    def safe_get(col):
        if col in current:
            return current[col]
        if col == "Engaged":
            return current.get("Engaged Visits", current.get("Engaged_Visits", 0))
        return 0

    metrics = {
        "Spend": safe_get("Spend"),
        "Traffic": safe_get("Traffic"),
        "Engaged": safe_get("Engaged"),
        "EVR": (
            safe_get("Engaged") / safe_get("Traffic")
            if safe_get("Traffic") > 0 else 0
        )
    }

    deltas = {}
    if prev is not None:
        prev_eng = prev.get("Engaged Visits", prev.get("Engaged_Visits", 0))
        prev_evr = prev_eng / prev["Traffic"] if prev["Traffic"] > 0 else None

        deltas = {
            "Spend": get_delta(metrics["Spend"], prev["Spend"]),
            "Traffic": get_delta(metrics["Traffic"], prev["Traffic"]),
            "Engaged": get_delta(metrics["Engaged"], prev_eng),
            "EVR": get_delta(metrics["EVR"], prev_evr),
        }
    else:
        deltas = {k: None for k in metrics}

    return metrics, deltas



def add_kpi_block(slide, metrics, deltas):
    """
    Adds KPI numbers + icons + arrows with proper colouring.
    """

    for key, coords in KPI_COORDS.items():

        value = metrics[key]
        delta = deltas[key]

        # Format text
        formatted_value = format_kpi_value(key, value)
        icon = KPI_ICONS[key]

        # Arrow
        if delta is None:
            arrow = "→"
        elif delta > 0.02:
            arrow = "↑"
        elif delta < -0.02:
            arrow = "↓"
        else:
            arrow = "→"

        box = slide.shapes.add_textbox(coords["x"], coords["y"], Cm(6), Cm(2))
        tf = box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = f"{icon}  {formatted_value}   {arrow}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = colour_for_delta(delta)



# ============================================================================
# → SUMMARY TEXT + PLATFORM TABLE
# ============================================================================

def build_campaign_summary(df):
    df = df.sort_values(["Year", "Quarter", "Week"])

    if len(df) < 2:
        return "No week-on-week comparison available yet."

    curr = df.tail(1).iloc[0]
    prev = df.tail(2).iloc[0]

    def describe(curr, prev, name):
        if prev == 0:
            return f"{name} has no baseline."
        diff = (curr - prev) / prev
        if diff > 0.05:
            return f"{name} increased WoW."
        elif diff < -0.05:
            return f"{name} decreased WoW."
        return f"{name} held steady."

    spend_msg = describe(curr["Spend"], prev["Spend"], "Spend")
    traff_msg = describe(curr["Traffic"], prev["Traffic"], "Traffic")

    curr_eng = curr.get("Engaged Visits", curr.get("Engaged_Visits", 0))
    prev_eng = prev.get("Engaged Visits", prev.get("Engaged_Visits", 0))
    eng_msg = describe(curr_eng, prev_eng, "Engaged Visits")

    curr_evr = curr_eng / curr["Traffic"] if curr["Traffic"] else 0
    prev_evr = prev_eng / prev["Traffic"] if prev["Traffic"] else 0
    evr_msg = describe(curr_evr, prev_evr, "EVR")

    # Top Platform
    plat_rank = (
        df.groupby("Platform")["Engaged Visits"]
        .sum()
        .sort_values(ascending=False)
    )
    top_platform = plat_rank.index[0] if len(plat_rank) else "No platform data"

    return (
        f"• {spend_msg}\n"
        f"• {traff_msg}\n"
        f"• {eng_msg}\n"
        f"• {evr_msg}\n"
        f"• Highest contribution from **{top_platform}**"
    )


def add_summary_text(slide, text):
    box = slide.shapes.add_textbox(
        SUMMARY_COORDS["x"],
        SUMMARY_COORDS["y"],
        SUMMARY_COORDS["w"],
        SUMMARY_COORDS["h"],
    )
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.name = "Arial"


def add_platform_ranking_table(slide, df):
    """
    Adds a platform ranking table, sorted by Engaged Visits.
    """

    if "Platform" not in df.columns:
        return

    d = (
        df.groupby("Platform")
        .agg(
            Spend=("Spend", "sum"),
            Engaged=("Engaged Visits", "sum"),
            Traffic=("Traffic", "sum"),
        )
        .reset_index()
    )

    d["EVR"] = d["Engaged"] / d["Traffic"].replace(0, None)
    d = d.sort_values("Engaged", ascending=False).head(10)

    # Insert table
    rows, cols = len(d) + 1, 4
    table_shape = slide.shapes.add_table(
        rows, cols,
        TABLE_COORDS["x"], TABLE_COORDS["y"],
        TABLE_COORDS["w"], TABLE_COORDS["h"]
    )
    table = table_shape.table

    headers = ["Platform", "Spend", "Engaged", "EVR"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True

    for i, (_, row) in enumerate(d.iterrows(), start=1):
        table.cell(i, 0).text = str(row["Platform"])
        table.cell(i, 1).text = f"£{row['Spend']:,.0f}"
        table.cell(i, 2).text = f"{row['Engaged']:,.0f}"
        table.cell(i, 3).text = f"{row['EVR']*100:.1f}%"



# ============================================================================
# → MAIN BUILDER
# ============================================================================

class CampaignPPTBuilder:

    def __init__(self, template_path: str, output_path: str):
        self.template = template_path
        self.output = output_path

    def build(self, df_raw):

        prs = Presentation(self.template)
        layout = prs.slide_layouts[0]

        campaigns = sorted(df_raw["Campaign"].dropna().unique())
        chart_temp = "ppt_temp"
        os.makedirs(chart_temp, exist_ok=True)

        summary_totals = []

        for camp in campaigns:
            df_c = df_raw[df_raw["Campaign"] == camp].copy()
            slide = prs.slides.add_slide(layout)

            # Title
            if slide.shapes.title:
                slide.shapes.title.text = camp

            # KPIs
            metrics, deltas = extract_latest_kpis(df_c)
            add_kpi_block(slide, metrics, deltas)

            # Summary
            summary = build_campaign_summary(df_c)
            add_summary_text(slide, summary)

            # Platform table
            add_platform_ranking_table(slide, df_c)

            # Charts
            for label, metric in [
                ("Spend", "Spend"),
                ("Traffic", "Traffic"),
                ("Engaged", "Engaged Visits"),
                ("EVR", "EVR")
            ]:
                png = os.path.join(chart_temp, f"{camp}_{label}.png")
                fig = latest_week_dimension_bar(df_c, "Platform", metric)
                fig.write_image(png, scale=3)

                c = COORDS[label]
                slide.shapes.add_picture(png, c["x"], c["y"], width=c["w"], height=c["h"])

            # For final summary slide
            summary_totals.append({
                "Campaign": camp,
                **metrics
            })

        # -----------------------------------------------------------
        # FINAL SUMMARY SLIDE
        # -----------------------------------------------------------
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Overall Campaign Summary"

        df_sum = pd.DataFrame(summary_totals)

        total_spend = df_sum["Spend"].sum()
        total_traf = df_sum["Traffic"].sum()
        total_eng = df_sum["Engaged"].sum()
        avg_evr = total_eng / total_traf if total_traf > 0 else 0

        summary_text = (
            f"• Total Spend: £{total_spend:,.0f}\n"
            f"• Total Traffic: {total_traf:,.0f}\n"
            f"• Total Engaged Visits: {total_eng:,.0f}\n"
            f"• Average EVR: {avg_evr*100:.1f}%\n"
        )

        box = slide.shapes.add_textbox(Cm(0.6), Cm(3), Cm(25), Cm(4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = summary_text
        p.font.size = Pt(20)
        p.font.name = "Arial"

        prs.save(self.output)
        print(f"✓ PowerPoint created: {self.output}")
