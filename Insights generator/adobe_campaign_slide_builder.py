import pandas as pd
from pathlib import Path

from src.reporting import compute_campaign_kpis, campaign_summary_text
from src.visuals import (
    latest_week_dimension_bar,
)
from pptx import Presentation
from pptx.util import Cm, Inches, Pt

# ---------------------------------------------------
# CONFIG
# ---------------------------------------------------
BACKGROUND_IMAGE = "background_gradient.png"  # your gradient image
INPUT_EXCEL = "data.xlsx"              # <-- CHANGE THIS TO YOUR FILE
OUTPUT_DIR = Path("campaign_slides")
OUTPUT_DIR.mkdir(exist_ok=True)

CHART_WIDTH = Inches(3.3)
CHART_HEIGHT = Inches(2.4)

# campaign slide positions (same for all)
POS = {
    "title_x": Cm(1),
    "title_y": Cm(1),

    "spend_x": Cm(0.41),
    "spend_y": Cm(12.48),

    "traffic_x": Cm(8.44),
    "traffic_y": Cm(12.48),

    "engaged_x": Cm(16.48),
    "engaged_y": Cm(12.48),

    "evr_x": Cm(24.71),
    "evr_y": Cm(12.48),
}

# ---------------------------------------------------
# LOAD DATA
# ---------------------------------------------------
print("📂 Loading dataset...")

df = pd.read_excel(INPUT_EXCEL)

required_cols = ["Campaign", "Year", "Quarter", "Week"]
if not set(required_cols).issubset(df.columns):
    raise ValueError(f"Input file is missing required columns: {required_cols}")

campaigns = sorted(df["Campaign"].dropna().unique())

print(f"🔍 Found {len(campaigns)} campaigns:")
for c in campaigns:
    print("  •", c)

# ---------------------------------------------------
# FUNCTION: build slide for one campaign
# ---------------------------------------------------
def build_campaign_slide(campaign_name, df_campaign):

    prs = Presentation()
    blank = prs.slide_layouts[6]  # empty slide
    slide = prs.slides.add_slide(blank)

    # background image
    slide.shapes.add_picture(BACKGROUND_IMAGE, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

    # ---- Title ----
    title_box = slide.shapes.add_textbox(POS["title_x"], POS["title_y"], Cm(20), Cm(2))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = campaign_name
    run.font.bold = True
    run.font.size = Pt(42)
    run.font.name = "Montserrat"

    # ---- KPI Summary ----
    kpis = compute_campaign_kpis(df_campaign)
    summary = campaign_summary_text(df_campaign)

    summary_box = slide.shapes.add_textbox(Cm(1), Cm(4), Cm(25), Cm(4))
    tf2 = summary_box.text_frame
    p = tf2.paragraphs[0]
    p.text = summary
    p.font.size = Pt(16)
    p.font.name = "Montserrat"

    # ---- CHARTS ----
    # Spend
    spend_fig = latest_week_dimension_bar(df_campaign, "Platform", "Spend", "Spend by Platform")
    spend_fig.write_image("_tmp_spend.png")
    slide.shapes.add_picture("_tmp_spend.png", POS["spend_x"], POS["spend_y"], width=CHART_WIDTH)

    # Traffic
    traf_fig = latest_week_dimension_bar(df_campaign, "Platform", "Traffic", "Traffic by Platform")
    traf_fig.write_image("_tmp_traf.png")
    slide.shapes.add_picture("_tmp_traf.png", POS["traffic_x"], POS["traffic_y"], width=CHART_WIDTH)

    # Engaged Visits
    ev_fig = latest_week_dimension_bar(df_campaign, "Platform", "Engaged Visits", "Engaged Visits by Platform")
    ev_fig.write_image("_tmp_eng.png")
    slide.shapes.add_picture("_tmp_eng.png", POS["engaged_x"], POS["engaged_y"], width=CHART_WIDTH)

    # EVR
    evr_fig = latest_week_dimension_bar(df_campaign, "Platform", "EVR", "EVR by Platform")
    evr_fig.write_image("_tmp_evr.png")
    slide.shapes.add_picture("_tmp_evr.png", POS["evr_x"], POS["evr_y"], width=CHART_WIDTH)

    return prs


# ---------------------------------------------------
# BUILD SLIDES FOR ALL CAMPAIGNS
# ---------------------------------------------------
print("\n📊 Generating slides...")

for c in campaigns:
    df_c = df[df["Campaign"] == c]

    prs = build_campaign_slide(c, df_c)
    outfile = OUTPUT_DIR / f"{c}.pptx"
    prs.save(outfile)

    print(f"  ✓ Saved: {outfile}")

print("\n✅ All campaign slides generated successfully!\n")
